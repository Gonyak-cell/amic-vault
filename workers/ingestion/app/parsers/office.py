from __future__ import annotations

import subprocess
import tempfile
from collections.abc import Iterable
from datetime import date, datetime, time
from decimal import Decimal
from io import BytesIO
from pathlib import Path

from app.converters.docx_to_pdf import legacy_office_signature, libreoffice_command

from .plaintext import extract_csv, extract_html, extract_plaintext
from .types import ExtractionResult

legacy_office_extensions = {"doc", "xls", "ppt"}
_legacy_targets = {
    "doc": ("txt", "txt"),
    "xls": ("csv", "csv"),
    "ppt": ("html", "html"),
}


def _stringify(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    if isinstance(value, date | time):
        return value.isoformat()
    if isinstance(value, Decimal):
        return format(value, "f")
    text = str(value).strip()
    return text


def _non_empty(values: Iterable[object]) -> list[str]:
    return [text for value in values if (text := _stringify(value))]


def _convert_legacy_office(payload: bytes, extension: str, timeout_seconds: int = 30) -> bytes:
    target = _legacy_targets.get(extension)
    if target is None:
        raise LegacyOfficeExtractionError("unsupported legacy office extension")
    if not payload.startswith(legacy_office_signature):
        raise LegacyOfficeExtractionError("input is not a legacy office compound payload")

    convert_to, output_extension = target
    with tempfile.TemporaryDirectory(prefix="amic-legacy-office-") as tmp:
        workdir = Path(tmp)
        source = workdir / f"source.{extension}"
        source.write_bytes(payload)
        try:
            subprocess.run(
                [
                    libreoffice_command(),
                    "--headless",
                    "--nologo",
                    "--nofirststartwizard",
                    "--convert-to",
                    convert_to,
                    "--outdir",
                    str(workdir),
                    str(source),
                ],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=timeout_seconds,
            )
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
            raise LegacyOfficeExtractionError("libreoffice text conversion failed") from exc

        output = workdir / f"source.{output_extension}"
        if not output.exists():
            raise LegacyOfficeExtractionError("libreoffice did not write a text derivative")
        return output.read_bytes()


def extract_legacy_office(payload: bytes, extension: str) -> ExtractionResult:
    if extension not in legacy_office_extensions:
        return ExtractionResult.failed("failed", "UNSUPPORTED_FILE_TYPE")
    try:
        text_payload = _convert_legacy_office(payload, extension)
    except LegacyOfficeExtractionError:
        return ExtractionResult.failed(extension, "LEGACY_OFFICE_TEXT_EXTRACTION_FAILED")

    if extension == "xls":
        result = extract_csv(text_payload)
    elif extension == "ppt":
        result = extract_html(text_payload)
    else:
        result = extract_plaintext(text_payload, extension)
    if result.status != "ready":
        return ExtractionResult.failed(extension, result.failure_reason_code or "LEGACY_OFFICE_TEXT_EMPTY")
    return ExtractionResult.ready(extension, result.body_text, result.confidence)


def extract_xlsx(payload: bytes) -> ExtractionResult:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(BytesIO(payload), read_only=True, data_only=True)
    except Exception:
        return ExtractionResult.failed("xlsx", "PROTECTED_XLSX_OR_INVALID")

    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows = [
                " | ".join(cells)
                for row in sheet.iter_rows(values_only=True)
                if (cells := _non_empty(row))
            ]
            if rows:
                parts.append(sheet.title)
                parts.extend(rows)
    finally:
        workbook.close()

    body_text = "\n".join(parts).strip()
    if not body_text:
        return ExtractionResult.failed("xlsx", "XLSX_TEXT_EMPTY")
    return ExtractionResult.ready("xlsx", body_text)


def _shape_text(shape: object) -> list[str]:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text_frame = getattr(shape, "text_frame")
        for paragraph in text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                parts.append(text)
    if getattr(shape, "has_table", False):
        table = getattr(shape, "table")
        for row in table.rows:
            cells = _non_empty(cell.text for cell in row.cells)
            if cells:
                parts.append(" | ".join(cells))
    for child in getattr(shape, "shapes", []):
        parts.extend(_shape_text(child))
    return parts


def extract_pptx(payload: bytes) -> ExtractionResult:
    try:
        from pptx import Presentation

        presentation = Presentation(BytesIO(payload))
    except Exception:
        return ExtractionResult.failed("pptx", "PROTECTED_PPTX_OR_INVALID")

    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            slide_parts.extend(_shape_text(shape))
        if slide_parts:
            parts.append(f"Slide {index}")
            parts.extend(slide_parts)

    body_text = "\n".join(parts).strip()
    if not body_text:
        return ExtractionResult.failed("pptx", "PPTX_TEXT_EMPTY")
    return ExtractionResult.ready("pptx", body_text)


class LegacyOfficeExtractionError(RuntimeError):
    pass
