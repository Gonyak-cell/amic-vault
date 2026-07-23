from __future__ import annotations

import subprocess
from datetime import datetime, timedelta, timezone
from hashlib import sha256
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from uuid import uuid4
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from reportlab.pdfgen import canvas

from app import extract_router
from app.converters.docx_to_pdf import libreoffice_command
from app.main import app
from app.storage_client import WorkerStoredObject

TENANT_ID = "11111111-1111-4111-8111-111111111111"
VERSION_ID = "11111111-1111-4111-8111-111111111155"
MATTER_ID = "11111111-1111-4111-8111-111111111122"
DOCUMENT_ID = "11111111-1111-4111-8111-111111111133"
FILE_OBJECT_ID = "11111111-1111-4111-8111-111111111144"

client = TestClient(app)
_stored_objects: dict[str, WorkerStoredObject] = {}


def _fake_read_ingestion_object(job):
    return _stored_objects[job.objectKey]


def _loopback_identity_headers() -> dict[str, str]:
    request_id = str(uuid4())
    nonce = str(uuid4())
    expires_at = (datetime.now(timezone.utc) + timedelta(minutes=3)).replace(microsecond=0)
    return {
        "x-amic-dev-loopback-identity": "true",
        "x-amic-request-id": request_id,
        "x-amic-ingestion-nonce": nonce,
        "x-amic-ingestion-expires-at": expires_at.strftime("%Y-%m-%dT%H:%M:%SZ"),
    }


def _content_type(filename: str) -> str:
    extension = filename.rsplit(".", 1)[-1].lower()
    return {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "doc": "application/msword",
        "xls": "application/vnd.ms-excel",
        "ppt": "application/vnd.ms-powerpoint",
        "hwp": "application/x-hwp",
        "hwpx": "application/hwp+zip",
        "txt": "text/plain",
        "md": "text/markdown",
        "csv": "text/csv",
        "html": "text/html",
        "htm": "text/html",
    }[extension]


def _post_extract(filename: str, payload: bytes, tenant_id: str = TENANT_ID):
    request_id = str(uuid4())
    nonce = str(uuid4())
    object_key = f"tenants/{tenant_id}/matters/{MATTER_ID}/documents/{DOCUMENT_ID}/{FILE_OBJECT_ID}"
    _stored_objects[object_key] = WorkerStoredObject(payload, _content_type(filename))
    app.state.ingestion_storage_reader = _fake_read_ingestion_object
    expires_at = (datetime.now(timezone.utc) + timedelta(minutes=3)).replace(microsecond=0)
    expires_at_value = expires_at.strftime("%Y-%m-%dT%H:%M:%SZ")
    return client.post(
        "/extract",
        json={
            "tenantId": tenant_id,
            "documentId": DOCUMENT_ID,
            "versionId": VERSION_ID,
            "fileObjectId": FILE_OBJECT_ID,
            "storageAlias": "primary",
            "objectKey": object_key,
            "objectVersion": "b" * 64,
            "sha256": sha256(payload).hexdigest(),
            "sizeBytes": len(payload),
            "parserProfile": "extract",
            "requestId": request_id,
            "expiresAt": expires_at_value,
        },
        headers={
            "x-amic-dev-loopback-identity": "true",
            "x-amic-request-id": request_id,
            "x-amic-ingestion-nonce": nonce,
            "x-amic-ingestion-expires-at": expires_at_value,
        },
    )


def _text_pdf(text: str) -> bytes:
    buffer = BytesIO()
    page = canvas.Canvas(buffer)
    page.drawString(72, 720, text)
    page.showPage()
    page.save()
    return buffer.getvalue()


def _blank_pdf() -> bytes:
    buffer = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.write(buffer)
    return buffer.getvalue()


def _encrypted_pdf() -> bytes:
    buffer = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.encrypt("fixture-password")
    writer.write(buffer)
    return buffer.getvalue()


def _docx_with_table_and_footnote() -> bytes:
    document = Document()
    document.add_paragraph("Matter summary paragraph")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Clause"
    table.rows[0].cells[1].text = "Status"
    buffer = BytesIO()
    document.save(buffer)
    source = BytesIO(buffer.getvalue())
    output = BytesIO()
    with ZipFile(source) as src, ZipFile(output, "w", ZIP_DEFLATED) as dst:
        for info in src.infolist():
            dst.writestr(info, src.read(info.filename))
        dst.writestr(
            "word/footnotes.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<w:footnotes xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:footnote w:id="2"><w:p><w:r><w:t>Footnote fixture text</w:t></w:r></w:p></w:footnote>
</w:footnotes>""",
        )
    return output.getvalue()


def _hwpx(section_texts: list[str]) -> bytes:
    output = BytesIO()
    with ZipFile(output, "w", ZIP_DEFLATED) as archive:
        archive.writestr("mimetype", "application/hwp+zip")
        for index, text in enumerate(section_texts):
            archive.writestr(
                f"Contents/section{index}.xml",
                f"""<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2016/section">
  <hp:p><hp:run><hp:t>{text}</hp:t></hp:run></hp:p>
</hp:sec>""",
            )
    return output.getvalue()


def _hwp_binary_fixture() -> bytes:
    return b"\xd0\xcf\x11\xe0" + b"\x00" * 64


def _xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Matter"
    sheet.append(["Clause", "Status"])
    sheet.append(["Spreadsheet covenant", "Open"])
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def _pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Deck summary"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1))
    text_box.text = "Presentation covenant"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _legacy_office(openxml_payload: bytes, source_extension: str, target_extension: str) -> bytes:
    with TemporaryDirectory(prefix="amic-test-legacy-office-") as tmp:
        workdir = Path(tmp)
        source = workdir / f"source.{source_extension}"
        source.write_bytes(openxml_payload)
        subprocess.run(
            [
                libreoffice_command(),
                "--headless",
                "--nologo",
                "--nofirststartwizard",
                "--convert-to",
                target_extension,
                "--outdir",
                str(workdir),
                str(source),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=30,
        )
        output = workdir / f"source.{target_extension}"
        assert output.exists(), f"LibreOffice did not create {output.name}"
        return output.read_bytes()


def test_pdf_text_layer_extraction_preserves_content() -> None:
    response = _post_extract("fixture.pdf", _text_pdf("PDF fixture first page"))
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["status"] == "ready"
    assert body["extraction_method"] == "pdf_text"
    assert body["confidence"] == 1.0
    assert "PDF fixture first page" in body["body_text"]


def test_storage_url_input_is_rejected_without_network_access() -> None:
    headers = _loopback_identity_headers()
    response = client.post(
        "/extract",
        json={"storage_url": "https://storage.local/presigned-fixture.pdf"},
        headers=headers,
    )
    assert response.status_code == 400
    assert "VALIDATION_FAILED" in response.text


def test_blank_pdf_is_ocr_pending_without_external_ocr() -> None:
    response = _post_extract("blank.pdf", _blank_pdf())
    assert response.status_code == 200, response.text
    assert response.json() == {
        "status": "ocr_pending",
        "extraction_method": "ocr_required",
        "body_text": "",
        "confidence": 0.0,
        "failure_reason_code": None,
    }


def test_encrypted_pdf_returns_explicit_failure_without_body_text() -> None:
    response = _post_extract("encrypted.pdf", _encrypted_pdf())
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["status"] == "failed"
    assert body["failure_reason_code"] == "ENCRYPTED_PDF"
    assert body["body_text"] == ""


def test_docx_extraction_includes_body_table_and_footnote_text() -> None:
    response = _post_extract("fixture.docx", _docx_with_table_and_footnote())
    assert response.status_code == 200, response.text
    body_text = response.json()["body_text"]
    assert "Matter summary paragraph" in body_text
    assert "Clause | Status" in body_text
    assert "Footnote fixture text" in body_text


def test_plaintext_markdown_csv_and_html_extract_searchable_body_text() -> None:
    cases = [
        (
            "fixture.txt",
            "계약서 본문 plain text".encode("utf-8"),
            "text",
            "계약서 본문 plain text",
        ),
        (
            "fixture.md",
            "# 검토 메모\n\n- Markdown clause token".encode("utf-8"),
            "markdown",
            "Markdown clause token",
        ),
        (
            "fixture.csv",
            "clause,status\n비밀유지,open\n".encode("utf-8"),
            "csv",
            "비밀유지 | open",
        ),
        (
            "fixture.html",
            b"<html><head><style>.x{}</style></head><body><h1>HTML Title</h1><p>Visible clause</p><script>hidden()</script></body></html>",
            "html",
            "Visible clause",
        ),
        (
            "fixture.htm",
            "<html><body><p>HTM visible text</p></body></html>".encode("utf-8"),
            "html",
            "HTM visible text",
        ),
    ]
    for filename, payload, method, expected in cases:
        response = _post_extract(filename, payload)
        assert response.status_code == 200, response.text
        body = response.json()
        assert body["status"] == "ready"
        assert body["extraction_method"] == method
        assert expected in body["body_text"]


def test_plaintext_binary_payload_fails_without_body_text() -> None:
    response = _post_extract("bad.txt", b"\x00\x01\x02")
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["status"] == "failed"
    assert body["extraction_method"] == "text"
    assert body["failure_reason_code"] == "TEXT_DECODE_FAILED"
    assert body["body_text"] == ""


def test_xlsx_and_pptx_extract_searchable_body_text() -> None:
    cases = [
        ("fixture.xlsx", _xlsx(), "xlsx", "Spreadsheet covenant | Open"),
        ("fixture.pptx", _pptx(), "pptx", "Presentation covenant"),
    ]
    for filename, payload, method, expected in cases:
        response = _post_extract(filename, payload)
        assert response.status_code == 200, response.text
        body = response.json()
        assert body["status"] == "ready"
        assert body["extraction_method"] == method
        assert expected in body["body_text"]


def test_legacy_doc_xls_and_ppt_extract_searchable_body_text() -> None:
    cases = [
        (
            "fixture.doc",
            _legacy_office(_docx_with_table_and_footnote(), "docx", "doc"),
            "doc",
            "Matter summary paragraph",
        ),
        (
            "fixture.xls",
            _legacy_office(_xlsx(), "xlsx", "xls"),
            "xls",
            "Spreadsheet covenant | Open",
        ),
        (
            "fixture.ppt",
            _legacy_office(_pptx(), "pptx", "ppt"),
            "ppt",
            "Deck summary",
        ),
    ]
    for filename, payload, method, expected in cases:
        response = _post_extract(filename, payload)
        assert response.status_code == 200, response.text
        body = response.json()
        assert body["status"] == "ready"
        assert body["extraction_method"] == method
        assert expected in body["body_text"]


def test_hwpx_extraction_fixtures_cover_five_deidentified_shapes() -> None:
    fixtures = {
        "basic.hwpx": ["Basic HWPX fixture"],
        "table.hwpx": ["Header A", "Cell B"],
        "image-heavy.hwpx": ["Caption only fixture text"],
        "large.hwpx": [f"Large paragraph {index}" for index in range(25)],
        "legacy.hwpx": ["Legacy style HWPX fixture"],
    }
    for filename, sections in fixtures.items():
        response = _post_extract(filename, _hwpx(sections))
        assert response.status_code == 200, response.text
        body = response.json()
        assert body["status"] == "ready"
        assert body["extraction_method"] == "hwpx"
        for expected in sections:
            assert expected in body["body_text"]


def test_hwpx_endpoint_rejects_hwp_binary_without_binary_parser() -> None:
    response = _post_extract("binary.hwpx", b"\xd0\xcf\x11\xe0" + b"not-real-document")
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["status"] == "failed"
    assert body["failure_reason_code"] == "UNSUPPORTED_HWP_BINARY"
    assert body["body_text"] == ""


def test_hwp5_binary_extraction_uses_hwp5txt(monkeypatch) -> None:
    def fake_run(command, *, profile_name, cwd, check, timeout_seconds):
        assert command[0] == "hwp5txt"
        assert command[1].endswith("source.hwp")
        assert profile_name == "extract"
        assert cwd
        assert check is False
        assert timeout_seconds == 30
        return subprocess.CompletedProcess(
            command,
            0,
            stdout="법원 제출 서면 청구취지 HWP 본문".encode("utf-8"),
            stderr=b"",
        )

    monkeypatch.setattr("app.parsers.hwp_binary.run_bounded_subprocess", fake_run)

    response = _post_extract("legacy.hwp", _hwp_binary_fixture())
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["status"] == "ready"
    assert body["extraction_method"] == "hwp5"
    assert "법원 제출 서면 청구취지" in body["body_text"]
    assert body["failure_reason_code"] is None


def test_hwp5_binary_failure_returns_explicit_reason(monkeypatch) -> None:
    def fake_run(command, **_):
        return subprocess.CompletedProcess(
            command,
            1,
            stdout=b"",
            stderr=b"encrypted distribution document requires password",
        )

    monkeypatch.setattr("app.parsers.hwp_binary.run_bounded_subprocess", fake_run)

    response = _post_extract("legacy.hwp", _hwp_binary_fixture())
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["status"] == "failed"
    assert body["extraction_method"] == "hwp5"
    assert body["failure_reason_code"] == "HWP_ENCRYPTED_OR_DRM"
    assert body["body_text"] == ""


def test_tenant_header_mismatch_fails_closed() -> None:
    headers = _loopback_identity_headers()
    response = client.post(
        "/extract",
        json={},
        headers={**headers, "x-amic-tenant-id": "22222222-2222-4222-8222-222222222222"},
    )
    assert response.status_code == 400
    assert "VALIDATION_FAILED" in response.text
